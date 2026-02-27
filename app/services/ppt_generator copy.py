from pptx import Presentation
import io
import requests


def replace_text_preserve_format(text_frame, replacements):
    for paragraph in text_frame.paragraphs:
        full_text = "".join(run.text for run in paragraph.runs)

        new_text = full_text
        for key, value in replacements.items():
            new_text = new_text.replace(key, value)

        if new_text != full_text:
            for run in paragraph.runs:
                run.text = ""

            if paragraph.runs:
                paragraph.runs[0].text = new_text
            else:
                paragraph.add_run().text = new_text


def generate_ppt(template_url: str, replacements: dict):

    response = requests.get(template_url)
    response.raise_for_status()

    ppt_stream = io.BytesIO(response.content)
    prs = Presentation(ppt_stream)

    for slide in prs.slides:
        for shape in slide.shapes:

            # 🔹 Texto normal
            if shape.has_text_frame:
                replace_text_preserve_format(shape.text_frame, replacements)

            # 🔹 Texto en tablas
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text_frame:
                            replace_text_preserve_format(cell.text_frame, replacements)

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)

    return output