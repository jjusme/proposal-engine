from pptx import Presentation
import io
import requests


def replace_text_preserve_format(text_frame, replacements):
    for paragraph in text_frame.paragraphs:
        full_text = "".join(run.text for run in paragraph.runs)

        new_text = full_text
        for key, value in replacements.items():
            new_text = new_text.replace(key, value)

        if new_text != full_text:
            for run in paragraph.runs:
                run.text = ""

            if paragraph.runs:
                paragraph.runs[0].text = new_text
            else:
                paragraph.add_run().text = new_text


def replace_logo(prs, logo_url):

    for slide in prs.slides:
        # usamos list() porque vamos a modificar shapes
        for shape in list(slide.shapes):

            if shape.name == "CLIENT_LOGO":

                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height

                # 🔴 Si no viene logo → eliminar placeholder
                if not logo_url:
                    slide.shapes._spTree.remove(shape._element)
                    continue

                # Descargar imagen
                response = requests.get(logo_url)
                response.raise_for_status()
                image_stream = io.BytesIO(response.content)

                # Eliminar placeholder
                slide.shapes._spTree.remove(shape._element)

                # Insertar fijando altura (mantiene proporción)
                picture = slide.shapes.add_picture(
                    image_stream,
                    left,
                    top,
                    height=height
                )

                # Centrar horizontalmente dentro del área original
                picture.left = left + (width - picture.width) // 2


def generate_ppt(template_url: str, replacements: dict, logo_url: str = None):

    response = requests.get(template_url)
    response.raise_for_status()

    ppt_stream = io.BytesIO(response.content)
    prs = Presentation(ppt_stream)

    for slide in prs.slides:
        for shape in slide.shapes:

            if shape.has_text_frame:
                replace_text_preserve_format(shape.text_frame, replacements)

            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text_frame:
                            replace_text_preserve_format(cell.text_frame, replacements)

    replace_logo(prs, logo_url)

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)

    return output